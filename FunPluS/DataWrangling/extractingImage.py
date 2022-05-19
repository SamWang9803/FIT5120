filename = 'Presentation3.pptx'
# the file name where the icons are extracted and saved

slide_folder = {0:'animal',
               1: 'bugs',
               2: 'ocean',
               3: 'house',
               4: 'food',
               5: 'nature',
               6: 'clothes',
               7: 'transport',
               8: 'sports',
               9: 'magic',
               10:'art',
               11:'body'}
# the slide index and the respective category of the icons


# https://stackoverflow.com/questions/52491656/extracting-images-from-presentation-file

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

'''
Parameters
------------
prs: the Presentation object
slideIdx: the index number of the slide for parsing

Returns
------------
shape: the shape file parsed from the given presentation slide

'''
def iter_picture_shapes(prs, slideIdx):
    for shape in prs.slides[slideIdx].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        # parse only the PICTURE type element from the presentation slide
            yield shape

for slideIdx, folderName in slide_folder.items():
    n = 0
    if not os.path.exists(folderName):
        os.makedirs(folderName)
        # create the respective directory of each icon category

    for picture in iter_picture_shapes(Presentation(filename), slideIdx):
        image = picture.image
        # ---get image "file" contents---
        image_bytes = image.blob
        # ---make up a name for the file, e.g. 'image.jpg'---
        image_filename = folderName+'/image'+str(n)+'.png'
        with open(image_filename, 'wb') as f:
            f.write(image_bytes)
        n+=1
